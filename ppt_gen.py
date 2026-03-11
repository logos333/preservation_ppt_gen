import shutil
import os
import copy
from datetime import datetime
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt, Cm

# ==========================================
# КОНФИГУРАЦИЯ
# ==========================================

MAX_IMAGES_PER_ROW = 3
MAX_ROWS_PER_SLIDE = 2
IMAGES_PER_SLIDE = MAX_IMAGES_PER_ROW * MAX_ROWS_PER_SLIDE  # 6

# ==========================================
# 1. ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ (Математика и базовые операции)
# ==========================================

def replace_text(shape, target_text, new_text):
    """Ищет target_text и меняет на new_text с сохранением стиля."""
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if target_text in run.text:
                run.text = run.text.replace(target_text, new_text)

def get_working_area(slide, slide_width):
    """Вычисляет доступную зону между самым верхним и самым нижним текстом."""
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    if len(text_shapes) < 2:
        return Cm(1), Cm(3), slide_width - Cm(2), Cm(13)

    top_shape = min(text_shapes, key=lambda s: s.top)
    bottom_shape = max(text_shapes, key=lambda s: s.top)

    box_top = top_shape.top + top_shape.height + Cm(0.5)
    box_bottom = bottom_shape.top - Cm(0.5)
    box_height = max(box_bottom - box_top, Cm(1))
    box_left = Cm(1)
    box_width = slide_width - (box_left * 2)

    return box_left, box_top, box_width, box_height

def calculate_grid_layout(image_paths, box_left, box_top, box_width, box_height, spacing_pt=15):
    """Рассчитывает координаты для компоновки картинок в сетку (макс. 3 в ряду, макс. 2 ряда)."""
    if not image_paths:
        return []

    spacing = Pt(spacing_pt)

    # Разбиваем на строки по MAX_IMAGES_PER_ROW
    rows = []
    for i in range(0, len(image_paths), MAX_IMAGES_PER_ROW):
        rows.append(image_paths[i:i + MAX_IMAGES_PER_ROW])

    num_rows = len(rows)
    row_spacing = spacing if num_rows > 1 else 0
    row_height = (box_height - (num_rows - 1) * row_spacing) / num_rows

    layout = []
    current_y = box_top

    for row_images in rows:
        # Считаем aspect ratio для каждой картинки в ряду
        aspect_ratios = []
        for path in row_images:
            with Image.open(path) as img:
                aspect_ratios.append(img.width / img.height)

        n = len(row_images)
        sum_a = sum(aspect_ratios)

        if sum_a == 0:
            current_y += row_height + row_spacing
            continue

        # Идеальная высота чтобы заполнить ширину
        ideal_h = (box_width - (n - 1) * spacing) / sum_a

        # Ограничиваем высотой строки
        if ideal_h > row_height:
            h = row_height
            total_w = h * sum_a + (n - 1) * spacing
            start_x = box_left + (box_width - total_w) / 2
        else:
            h = ideal_h
            start_x = box_left

        # Центрируем по вертикали внутри строки
        y = current_y + (row_height - h) / 2

        current_x = start_x
        for i, a in enumerate(aspect_ratios):
            w = h * a
            layout.append({
                'path': row_images[i],
                'left': round(current_x), 'top': round(y),
                'width': round(w), 'height': round(h)
            })
            current_x += w + spacing

        current_y += row_height + row_spacing

    return layout

# ==========================================
# 2. БИЗНЕС-ЛОГИКА (Работа с папками и слайдами)
# ==========================================

def build_image_dictionary(images_folder):
    """Сканирует папку с фото и возвращает словарь сгруппированных картинок по тегам."""
    image_dict = {}
    if not os.path.exists(images_folder):
        print(f"Внимание: Папка '{images_folder}' не найдена!")
        return image_dict

    for filename in sorted(os.listdir(images_folder)):
        if filename.lower().endswith(('.png', '.jpg', '.jpeg')):
            name_without_ext = os.path.splitext(filename)[0]
            parts = name_without_ext.split('-')
            if name_without_ext.startswith('101-') or name_without_ext.startswith('201-'):
                img_tag = parts[1].strip() if len(parts) > 1 else parts[0].strip()
            else:
                img_tag = parts[0].strip()
            full_path = os.path.join(images_folder, filename)
            image_dict.setdefault(img_tag, []).append(full_path)
            
    return image_dict

def get_slide_tag(slide):
    """Находит верхний текст на слайде и извлекает из него тег (например, 'D903B')."""
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    if not text_shapes:
        return None
        
    top_shape = min(text_shapes, key=lambda s: s.top)
    top_text = top_shape.text.strip().split('\n')[0] 
    
    if '-' in top_text:
        return top_text.split('-')[1].strip()
    return None

def clear_slide_images(slide):
    """Удаляет все старые картинки с переданного слайда."""
    images_to_delete = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    for shape in images_to_delete:
        sp = shape._element
        sp.getparent().remove(sp)

def duplicate_slide(prs, slide_index):
    """Дублирует слайд по индексу и вставляет копию сразу после оригинала. Возвращает новый слайд."""
    source_slide = prs.slides[slide_index]
    slide_layout = source_slide.slide_layout

    # Создаём новый слайд с тем же layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Копируем все элементы исходного слайда через XML deep copy
    for shape in source_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

        # Копируем связанные ресурсы (картинки, медиа)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blip = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
            if blip is not None:
                embed_rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if embed_rId:
                    source_part = source_slide.part
                    target_part = new_slide.part
                    rel = source_part.rels[embed_rId]
                    new_rId = target_part.relate_to(rel.target_part, rel.reltype)
                    # Обновляем rId в скопированном элементе
                    new_blip = el.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                    if new_blip is not None:
                        new_blip.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed', new_rId)

    # Убираем placeholder-элементы, которые добавились из layout
    for ph_shape in list(new_slide.placeholders):
        sp = ph_shape._element
        sp.getparent().remove(sp)

    # Перемещаем слайд сразу после оригинала
    slide_list = prs.slides._sldIdLst
    slide_ids = list(slide_list)
    # Последний добавленный — наш новый слайд, перемещаем на позицию slide_index + 1
    new_slide_id = slide_ids[-1]
    slide_list.remove(new_slide_id)
    slide_list.insert(slide_index + 1, new_slide_id)

    return new_slide

def append_part_label(slide, part_number):
    """Добавляет ' (part N)' к первой строке верхнего текста на слайде."""
    text_shapes = [s for s in slide.shapes if s.has_text_frame and s.text.strip()]
    if not text_shapes:
        return

    top_shape = min(text_shapes, key=lambda s: s.top)
    first_paragraph = top_shape.text_frame.paragraphs[0]
    
    if first_paragraph.runs:
        last_run = first_paragraph.runs[-1]
        last_run.text += f" (part {part_number})"

def place_images_on_slide(slide, image_paths, slide_width):
    """Очищает старые картинки и размещает новые в сетке."""
    clear_slide_images(slide)
    b_left, b_top, b_width, b_height = get_working_area(slide, slide_width)
    layouts = calculate_grid_layout(image_paths, b_left, b_top, b_width, b_height)

    for item in layouts:
        slide.shapes.add_picture(
            item['path'],
            item['left'],
            item['top'],
            item['width'],
            item['height']
        )

def process_slide(prs, slide_index, current_date, image_dict, slide_width):
    """Главная функция обработки одного слайда. Возвращает кол-во добавленных слайдов."""
    slide = prs.slides[slide_index]

    # 1. Замена даты на всех слайдах без исключения
    for shape in slide.shapes:
        replace_text(shape, '{{DATE}}', current_date)
        
    # 2. Логика замены картинок работает только с 3-й страницы (индекс 2)
    if slide_index < 2:
        return 0

    # 3. Ищем тег. Если его нет — картинки на этом слайде не трогаем
    slide_tag = get_slide_tag(slide)
    if not slide_tag:
        return 0

    # 4. Ищем новые картинки для этого тега. Если их нет — ничего не делаем
    matched_images = sorted(image_dict.get(slide_tag, []))
    if not matched_images:
        return 0

    # 5. Разбиваем на группы по IMAGES_PER_SLIDE
    chunks = [matched_images[i:i + IMAGES_PER_SLIDE]
              for i in range(0, len(matched_images), IMAGES_PER_SLIDE)]

    # 6. Первая порция — на оригинальный слайд
    place_images_on_slide(slide, chunks[0], slide_width)

    # 7. Остальные порции — на дублированные слайды
    added_slides = 0
    for part_idx, chunk in enumerate(chunks[1:], start=2):
        new_slide = duplicate_slide(prs, slide_index + added_slides)
        added_slides += 1
        append_part_label(new_slide, part_idx)
        place_images_on_slide(new_slide, chunk, slide_width)

    if len(chunks) > 1:
        # Добавляем "(part 1)" к оригинальному слайду если есть продолжение
        append_part_label(slide, 1)

    return added_slides

# ==========================================
# 3. ТОЧКА ВХОДА
# ==========================================

def generate_presentation(
    images_folder: str = 'photos',
    template_path: str = 'template.pptx',
) -> str:
    """Генерирует презентацию из шаблона и фото. Возвращает путь к файлу."""
    current_date = datetime.now().strftime("%d-%B-%Y")
    new_file_path = f'MTO Preservation for {current_date}-Temur Khoshimov.pptx'

    print("Запуск скрипта...")

    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Файл шаблона '{template_path}' не найден!")

    # Подготовка
    shutil.copyfile(template_path, new_file_path)
    prs = Presentation(new_file_path)
    slide_width = prs.slide_width

    # Сбор данных
    image_dict = build_image_dictionary(images_folder)
    print(f"Найдено тегов в папке с фото: {len(image_dict)}")

    # Обработка слайдов (while-цикл, т.к. мы можем вставлять новые слайды)
    i = 0
    total_slides = len(prs.slides)
    while i < total_slides:
        added = process_slide(prs, i, current_date, image_dict, slide_width)
        i += 1 + added
        total_slides += added

    # Сохранение результата
    prs.save(new_file_path)
    print(f"Успешно! Презентация сохранена как '{new_file_path}'")
    return new_file_path

# Конструкция для правильного запуска Python-скриптов
if __name__ == '__main__':
    generate_presentation()